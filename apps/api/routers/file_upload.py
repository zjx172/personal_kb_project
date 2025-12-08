"""
文件上传相关路由（支持 PDF, DOCX, PPT, PPTX, Markdown）
"""
import asyncio
import logging
import uuid
from pathlib import Path
from typing import Optional

from fastapi import APIRouter, Depends, HTTPException, UploadFile, File, Form
from sqlalchemy.orm import Session
from datetime import datetime

from db import get_db, SessionLocal
from models import User, MarkdownDoc, KnowledgeBase
from auth import get_current_user
from schemas import UploadPdfResponse
from tasks import update_task, TaskStatus
from services.doc_service import upsert_markdown_doc_to_vectorstore
from services.storage import storage
from routers.pdf import get_or_create_default_knowledge_base

router = APIRouter(tags=["file-upload"])
logger = logging.getLogger(__name__)


def get_file_type(filename: str) -> str:
    """根据文件名获取文件类型"""
    ext = Path(filename).suffix.lower()
    if ext == ".pdf":
        return "pdf"
    elif ext in [".docx"]:
        return "docx"
    elif ext in [".ppt", ".pptx"]:
        return "ppt"
    elif ext in [".md", ".markdown"]:
        return "markdown"
    else:
        return "unknown"


async def extract_file_text(file_path: Path, file_type: str) -> str:
    """提取文件文本内容"""
    if file_type == "pdf":
        # PDF 提取逻辑（复用现有逻辑）
        from langchain_community.document_loaders import PyPDFLoader
        try:
            loader = PyPDFLoader(str(file_path))
            docs = loader.load()
            return "\n\n".join([doc.page_content for doc in docs])
        except Exception:
            import pypdf
            text_parts = []
            with open(file_path, "rb") as f:
                pdf_reader = pypdf.PdfReader(f)
                for page in pdf_reader.pages:
                    try:
                        text_parts.append(page.extract_text())
                    except Exception:
                        pass
            return "\n\n".join(text_parts)
    
    elif file_type == "docx":
        # DOCX 提取
        try:
            from docx import Document
            doc = Document(file_path)
            paragraphs = [para.text for para in doc.paragraphs]
            return "\n\n".join(paragraphs)
        except ImportError:
            raise HTTPException(
                status_code=500,
                detail="python-docx 库未安装，无法处理 DOCX 文件"
            )
        except Exception as e:
            logger.error(f"提取 DOCX 文本失败: {e}")
            raise HTTPException(status_code=500, detail=f"提取 DOCX 文本失败: {str(e)}")
    
    elif file_type == "ppt":
        # PPT/PPTX 提取
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            text_parts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        text_parts.append(shape.text)
            return "\n\n".join(text_parts)
        except ImportError:
            raise HTTPException(
                status_code=500,
                detail="python-pptx 库未安装，无法处理 PPT 文件"
            )
        except Exception as e:
            logger.error(f"提取 PPT 文本失败: {e}")
            raise HTTPException(status_code=500, detail=f"提取 PPT 文本失败: {str(e)}")
    
    elif file_type == "markdown":
        # Markdown 文件直接读取
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return f.read()
        except Exception as e:
            logger.error(f"读取 Markdown 文件失败: {e}")
            raise HTTPException(status_code=500, detail=f"读取 Markdown 文件失败: {str(e)}")
    
    else:
        raise HTTPException(status_code=400, detail=f"不支持的文件类型: {file_type}")


async def process_file_extraction(
    task_id: str,
    file_path: Optional[Path] = None,
    file_url: Optional[str] = None,
    doc_id: str = "",
    user_id: str = "",
    knowledge_base_id: Optional[str] = None,
    title: Optional[str] = None,
    filename: str = "",
):
    """后台任务：处理文件提取"""
    try:
        file_type = get_file_type(filename)
        await update_task(task_id, TaskStatus.PROCESSING, 10, f"开始提取{file_type.upper()}文件文本内容...")

        # 如果使用 OSS，需要先下载文件到临时位置
        temp_file_path = None
        if file_url and not file_path:
            from services.storage import storage
            import tempfile
            ext = Path(filename).suffix
            temp_file_path = Path(tempfile.mktemp(suffix=ext))
            file_content = storage.download_file(file_url.split("/")[-1] if "/" in file_url else file_url)
            with open(temp_file_path, "wb") as f:
                f.write(file_content)
            file_path = temp_file_path
            logger.info(f"已从 OSS 下载文件到临时位置: {temp_file_path}")

        # 提取文件文本内容
        await update_task(task_id, TaskStatus.PROCESSING, 30, "正在提取文本内容...")
        file_text = await extract_file_text(file_path, file_type)

        if not file_text or not file_text.strip():
            raise Exception("未能从文件中提取到文本内容")

        await update_task(task_id, TaskStatus.PROCESSING, 70, "文本提取完成，正在保存文档...")

        # 保存到数据库
        now = datetime.utcnow()
        db = SessionLocal()
        try:
            # 确定知识库ID
            final_kb_id = knowledge_base_id
            if not final_kb_id:
                default_kb = get_or_create_default_knowledge_base(db, user_id)
                final_kb_id = default_kb.id
            else:
                kb = (
                    db.query(KnowledgeBase)
                    .filter(
                        KnowledgeBase.id == final_kb_id,
                        KnowledgeBase.user_id == user_id,
                    )
                    .first()
                )
                if not kb:
                    default_kb = get_or_create_default_knowledge_base(db, user_id)
                    final_kb_id = default_kb.id

            doc = MarkdownDoc(
                id=doc_id,
                user_id=user_id,
                knowledge_base_id=final_kb_id,
                title=title or filename,
                content=file_text,
                doc_type=file_type,
                pdf_file_path=str(file_path) if file_path else None,
                created_at=now,
                updated_at=now,
            )
            db.add(doc)
            db.commit()
            db.refresh(doc)
            logger.info(f"文档已保存到数据库，ID: {doc_id}")

            await update_task(
                task_id, TaskStatus.PROCESSING, 85, "文档已保存，正在同步到向量库..."
            )

            # 同步到向量库
            loop = asyncio.get_event_loop()
            await loop.run_in_executor(
                None,
                lambda: upsert_markdown_doc_to_vectorstore(doc),
            )

            await update_task(
                task_id, TaskStatus.COMPLETED, 100, "文件处理完成！", result=doc
            )
        finally:
            db.close()
            # 清理临时文件
            if temp_file_path and temp_file_path.exists():
                try:
                    temp_file_path.unlink()
                except Exception:
                    pass

    except Exception as e:
        logger.error(f"文件处理失败: {str(e)}", exc_info=True)
        await update_task(task_id, TaskStatus.FAILED, 0, f"处理失败: {str(e)}", error=str(e))


@router.post("/upload-file", response_model=UploadPdfResponse)
async def upload_file(
    file: UploadFile = File(...),
    title: Optional[str] = Form(None),
    knowledge_base_id: Optional[str] = Form(None),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user),
):
    """上传文件（支持 PDF, DOCX, PPT, PPTX, Markdown）并启动后台提取任务"""
    # 验证文件类型
    if not file.filename:
        raise HTTPException(status_code=400, detail="文件名不能为空")
    
    file_type = get_file_type(file.filename)
    if file_type == "unknown":
        raise HTTPException(
            status_code=400,
            detail="不支持的文件类型。支持的类型：PDF, DOCX, PPT, PPTX, Markdown"
        )

    logger.info(f"开始上传文件: {file.filename}, 类型: {file_type}, 用户: {current_user.id}")

    file_path = None
    file_url = None
    try:
        # 生成唯一文件名和任务ID
        doc_id = str(uuid.uuid4())
        task_id = str(uuid.uuid4())
        file_extension = Path(file.filename).suffix
        file_filename = f"{doc_id}{file_extension}"
        
        # 使用存储适配器上传文件
        storage_path = f"files/{file_filename}"

        logger.info(f"上传文件: {storage_path}")

        # 初始化任务
        await update_task(task_id, TaskStatus.PENDING, 0, "准备上传文件...")

        # 读取文件内容
        file.file.seek(0)
        file_content = await file.read()
        
        # 确定 content type
        content_type_map = {
            "pdf": "application/pdf",
            "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "ppt": "application/vnd.ms-powerpoint",
            "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "markdown": "text/markdown",
        }
        content_type = content_type_map.get(file_type, "application/octet-stream")
        
        # 上传到存储
        file_url = storage.upload_file(
            file_content=file_content,
            file_path=storage_path,
            content_type=content_type,
        )

        # 如果是本地存储，保存路径用于后续处理
        if hasattr(storage, 'base_dir'):
            file_path = Path(file_url)
        else:
            file_path = None
        
        file_size = len(file_content)
        logger.info(f"文件已上传，大小: {file_size} 字节，URL: {file_url}")

        await update_task(task_id, TaskStatus.PROCESSING, 5, "文件上传完成，开始处理...")

        # 启动后台任务
        asyncio.create_task(
            process_file_extraction(
                task_id=task_id,
                file_path=file_path,
                file_url=file_url,
                doc_id=doc_id,
                user_id=current_user.id,
                knowledge_base_id=knowledge_base_id,
                title=title,
                filename=file.filename,
            )
        )

        return UploadPdfResponse(
            task_id=task_id, message="文件已上传，正在后台处理..."
        )

    except HTTPException:
        raise
    except Exception as e:
        logger.error(f"文件上传失败: {str(e)}", exc_info=True)
        if file_path and file_path.exists():
            try:
                file_path.unlink()
            except Exception:
                pass
        raise HTTPException(status_code=500, detail=f"上传文件失败: {str(e)}")

